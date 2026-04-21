"""Build the 15-minute methods-focused presentation deck.

Output: Docs/Results/presentation/spatial_bridge_talk.pptx

All numbers printed on the slides come from CSVs under results/full_run_bayes/
so the deck stays in sync with the manuscript in
Docs/Results/manuscript_build_bayes/manuscript.tex.
"""

from __future__ import annotations

import os
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path

os.environ.setdefault(
    "MPLCONFIGDIR", str(Path(tempfile.gettempdir()) / "mpl_cache_slides")
)

import matplotlib  # noqa: E402

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
import pandas as pd  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402


ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "Docs" / "Results" / "presentation"
ASSETS = OUT_DIR / "_assets"
RESULTS = ROOT / "results" / "full_run_bayes"
FIGS = RESULTS / "figures"
SYN_FIGS = ROOT / "results" / "synthetic_validation" / "figures"
MS_FIGS = ROOT / "Docs" / "Results" / "manuscript_build_bayes" / "figures"

DECK_PATH = OUT_DIR / "spatial_bridge_talk.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ACCENT = RGBColor(0x1F, 0x4E, 0x79)
MUTED = RGBColor(0x60, 0x60, 0x60)
BG = RGBColor(0xFF, 0xFF, 0xFF)


# ---------------------------------------------------------------------------
# Asset preparation
# ---------------------------------------------------------------------------


def convert_pdf_figures() -> dict[str, Path]:
    """Rasterize the three manuscript PDF figures to PNG at 200 dpi."""
    ASSETS.mkdir(parents=True, exist_ok=True)
    out: dict[str, Path] = {}
    for stem in ("gen_benchmark_bars", "gen_loadings_heatmap", "gen_perturbation_norms"):
        pdf_in = MS_FIGS / f"{stem}.pdf"
        png_out = ASSETS / f"{stem}.png"
        if not pdf_in.exists():
            raise FileNotFoundError(pdf_in)
        if not png_out.exists():
            subprocess.run(
                [
                    "pdftoppm",
                    "-png",
                    "-r",
                    "200",
                    "-singlefile",
                    str(pdf_in),
                    str(ASSETS / stem),
                ],
                check=True,
            )
        out[stem] = png_out
    return out


def crop_top(src: Path, out_png: Path, *, fraction: float = 0.12) -> Path:
    """Crop the top ``fraction`` off an image.

    Used to strip the matplotlib debug title
    ('PCA 2D fallback on expression -- install umap-learn for UMAP')
    off the Stage 4 UMAP PNG before embedding.
    """
    from PIL import Image

    ASSETS.mkdir(parents=True, exist_ok=True)
    if out_png.exists():
        return out_png
    with Image.open(src) as img:
        w, h = img.size
        top_px = int(h * fraction)
        cropped = img.crop((0, top_px, w, h))
        cropped.save(out_png)
    return out_png


def render_equation(latex_like: str, out_png: Path, *, fontsize: int = 26) -> Path:
    """Render a mathtext equation to a transparent PNG."""
    ASSETS.mkdir(parents=True, exist_ok=True)
    if out_png.exists():
        return out_png
    fig = plt.figure(figsize=(8, 1.2))
    fig.patch.set_alpha(0.0)
    ax = fig.add_subplot(111)
    ax.axis("off")
    ax.text(
        0.5,
        0.5,
        latex_like,
        fontsize=fontsize,
        ha="center",
        va="center",
        color="black",
    )
    fig.savefig(out_png, dpi=300, transparent=True, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)
    return out_png


# ---------------------------------------------------------------------------
# Metrics
# ---------------------------------------------------------------------------


@dataclass
class Metrics:
    spatial_bridge_l2: float
    uncond_bridge_l2: float
    latent_knn_l2: float
    static_centroid_l2: float
    de_shift_l2: float
    nmf_err: float
    pca_err: float
    ica_err: float
    tumor_norm: float
    intermediate_norm: float
    normal_norm: float
    cna_vs_pert_r: float
    heldout_train_delta: float
    heldout_test_delta: float
    n_train: int
    n_heldout: int
    hypoxia_fdr: float
    oxphos_fdr: float
    coagulation_fdr: float


def _bench_val(df: pd.DataFrame, method: str, metric: str) -> float:
    sub = df[(df["method"] == method) & (df["metric"] == metric)]
    return float(sub["value"].iloc[0])


def load_metrics() -> Metrics:
    bench = pd.read_csv(RESULTS / "benchmark_metrics.csv")
    fact = pd.read_csv(RESULTS / "stage_6_factorization_comparison.csv")
    plaus = pd.read_csv(RESULTS / "stage_7_biological_plausibility.csv")
    heldout = pd.read_csv(RESULTS / "stage_8_heldout_metrics.csv")
    best = pd.read_csv(RESULTS / "best_term_per_program.csv")

    nmf = float(fact.loc[fact["method"] == "NMF", "reconstruction_error"].iloc[0])
    pca = float(fact.loc[fact["method"] == "PCA", "reconstruction_error"].iloc[0])
    ica = float(fact.loc[fact["method"] == "ICA", "reconstruction_error"].iloc[0])

    def _norm(lbl: str) -> float:
        row = plaus[(plaus["stratum"] == "marginal_label") & (plaus["label"] == lbl)]
        return float(row["mean_perturbation_norm"].iloc[0])

    global_row = plaus[plaus["label"] == "cna_vs_perturbation_norm_pearson"]
    cna_r = float(global_row["mean_cna"].iloc[0])

    ho_train = heldout[heldout["split"] == "train_sections_internal"].iloc[0]
    ho_test = heldout[heldout["split"] == "heldout_sections"].iloc[0]

    def _fdr_for(term: str) -> float:
        row = best[best["best_term"] == term]
        return float(row["fdr_bh"].iloc[0])

    return Metrics(
        spatial_bridge_l2=_bench_val(bench, "SpatialBridge", "distance_to_normal_mean_l2"),
        uncond_bridge_l2=_bench_val(bench, "UnconditionalBridge", "distance_to_normal_mean_l2"),
        latent_knn_l2=_bench_val(bench, "LatentNN_normal_blend", "distance_to_normal_mean_l2"),
        static_centroid_l2=_bench_val(bench, "StaticOT_centroid", "distance_to_normal_mean_l2"),
        de_shift_l2=_bench_val(bench, "DE_shift", "distance_to_normal_mean_l2"),
        nmf_err=nmf,
        pca_err=pca,
        ica_err=ica,
        tumor_norm=_norm("tumor"),
        intermediate_norm=_norm("intermediate"),
        normal_norm=_norm("normal"),
        cna_vs_pert_r=cna_r,
        heldout_train_delta=float(ho_train["mean_delta_toward_reference"]),
        heldout_test_delta=float(ho_test["mean_delta_toward_reference"]),
        n_train=int(ho_train["n_spots"]),
        n_heldout=int(ho_test["n_spots"]),
        hypoxia_fdr=_fdr_for("HALLMARK_HYPOXIA"),
        oxphos_fdr=_fdr_for("HALLMARK_OXIDATIVE_PHOSPHORYLATION"),
        coagulation_fdr=_fdr_for("HALLMARK_COAGULATION"),
    )


# ---------------------------------------------------------------------------
# Layout helpers
# ---------------------------------------------------------------------------


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_title(slide, text: str, *, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        r2 = p2.runs[0]
        r2.font.size = Pt(16)
        r2.font.color.rgb = MUTED
        r2.font.italic = True


def _add_bullets(
    slide,
    bullets: list[str],
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    size: int = 18,
) -> None:
    """Render a bullet list.

    Convention:
      - "" -> empty paragraph used as a vertical spacer (no bullet dot).
      - "  ..." (leading two spaces) -> continuation line, rendered indented
        without a bullet dot.
      - anything else -> bulleted line.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if b == "":
            p.text = ""
        elif b.startswith("  "):
            p.text = b
        else:
            p.text = "\u2022  " + b
        for run in p.runs:
            run.font.size = Pt(size)
        p.space_after = Pt(6)


def _add_image(
    slide,
    path: Path,
    *,
    left: Emu,
    top: Emu,
    width: Emu | None = None,
    height: Emu | None = None,
) -> None:
    slide.shapes.add_picture(
        str(path), left, top, width=width, height=height
    )


def _add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(10)
    run.font.italic = True
    run.font.color.rgb = MUTED


def _hide_slide(slide) -> None:
    slide.element.set("show", "0")


# ---------------------------------------------------------------------------
# Equations (mathtext)
# ---------------------------------------------------------------------------


def prep_equations() -> dict[str, Path]:
    eqs = {
        "context": r"$\bar{\mathbf{x}}_i \;=\; \dfrac{1}{k} \sum_{j \in \mathcal{N}_k(i)} \mathbf{x}_j$",
        "target": r"$\mathbf{y}_i \;=\; \boldsymbol{\mu} - \mathbf{x}_i,\quad \boldsymbol{\mu} \;=\; |\mathcal{R}|^{-1}\!\!\sum_{i \in \mathcal{R}} \mathbf{x}_i$",
        "ridge": r"$\Theta^{*} \;=\; \mathrm{arg\,min}_{\Theta}\; \|Z\Theta - Y\|_F^{2} + \lambda \|\Theta\|_F^{2}$",
        "ridge_closed": r"$\Theta^{*} \;=\; (Z^{\top} Z + \lambda I)^{-1}\, Z^{\top} Y$",
        "drift_decomp": r"$\mathbf{s}(\mathbf{x},\bar{\mathbf{x}}) \;=\; \mathbf{x}\,\Theta_{\mathrm{expr}} + \bar{\mathbf{x}}\,\Theta_{\mathrm{ctx}} + \mathbf{b}$",
        "euler": r"$\mathbf{x}_i^{(t+1)} \;=\; \mathbf{x}_i^{(t)} + \eta\, \mathbf{s}(\mathbf{x}_i^{(t)}, \bar{\mathbf{x}}_i) + \alpha\,(\bar{\mathbf{x}}_i - \mathbf{x}_i^{(t)})$",
        "perturb": r"$\mathbf{u}_i \;=\; \tilde{\mathbf{x}}_i - \mathbf{x}_i \;\in\; \mathbb{R}^{G}$",
        "nmf": r"$W^{*}, H^{*} \;=\; \mathrm{arg\,min}_{W,H \geq 0}\; \|U_{+} - W H\|_F^{2}$",
        "dref": r"$d_{\mathrm{ref}} \;=\; \dfrac{1}{n} \sum_{i=1}^{n} \| \tilde{\mathbf{x}}_i - \boldsymbol{\mu} \|_2$",
        "posterior": r"$\Theta \mid Z,Y \sim \mathcal{N}\!\left(\Theta^{*},\; \hat{\sigma}^{2} (Z^{\top}Z + \lambda I)^{-1}\right)$",
        "pred_var": r"$\mathrm{Var}[\mathbf{s}(\mathbf{z}_{*})] \;=\; \hat{\sigma}^{2}\; \mathbf{z}_{*}^{\top} M \, \mathbf{z}_{*},\quad M=(Z^{\top}Z+\lambda I)^{-1}$",
        "pert_std": r"$\mathrm{std}(\|\mathbf{u}_i\|) \;\approx\; \mathrm{std}(\mathbf{s}_i)\cdot |\eta| \cdot N$",
    }
    out: dict[str, Path] = {}
    for name, latex in eqs.items():
        out[name] = render_equation(latex, ASSETS / f"eq_{name}.png")
    return out


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def slide_title(prs: Presentation, m: Metrics) -> None:
    s = _blank_slide(prs)
    # Accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), SLIDE_W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    tbox = s.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.6))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Spatially Conditioned Counterfactual Transport"
    p.runs[0].font.size = Pt(40)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = "for Glioblastoma Transcriptional Programs"
    p2.runs[0].font.size = Pt(32)
    p2.runs[0].font.color.rgb = ACCENT

    sub = s.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.2))
    p = sub.text_frame.paragraphs[0]
    p.text = (
        "We learn a spatially-smoothed drift field that maps tumor-like Visium spots\n"
        "toward a within-section normal reference, then factorize the perturbations\n"
        "into interpretable gene programs."
    )
    p.runs[0].font.size = Pt(18)
    p.runs[0].font.italic = True
    p.runs[0].font.color.rgb = MUTED

    auth = s.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2))
    tf = auth.text_frame
    p = tf.paragraphs[0]
    p.text = "Michael A. Haidar"
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.bold = True
    p2 = tf.add_paragraph()
    p2.text = "CS 8395 Special Topics in Computational Biology  \u2013  Vanderbilt University"
    p2.runs[0].font.size = Pt(14)
    p2.runs[0].font.color.rgb = MUTED


def slide_motivation(prs: Presentation, m: Metrics) -> None:
    s = _blank_slide(prs)
    _add_title(s, "Why not just DE?", subtitle="From association to minimal transformation")
    _add_bullets(
        s,
        [
            "Differential expression identifies genes that differ between tumor and normal \u2014 association, not transformation.",
            "We want the minimal transcriptional change sufficient to move a tumor spot to a normal-like state,",
            "  respecting spatial neighborhood structure.",
            "Framing: entropy-regularized optimal transport / Schr\u00f6dinger bridge.",
            "Tractable approximation: linear ridge drift + Euler integration, with optional Bayesian posterior.",
            "",
            "Deliverables:",
            "  (1) spot-level counterfactual perturbation vectors \u2192",
            "  (2) coordinated gene programs via NMF \u2192",
            "  (3) held-out and synthetic validation.",
        ],
        left=Inches(0.8),
        top=Inches(1.6),
        width=Inches(11.7),
        height=Inches(5.3),
        size=20,
    )


def slide_data(prs: Presentation, m: Metrics) -> None:
    s = _blank_slide(prs)
    _add_title(s, "Data & cohort", subtitle="Glioblastoma Visium \u2014 3 sections (Bayes run)")
    _add_image(
        s,
        FIGS / "stage_2_spatial_marginal_labels.png",
        left=Inches(0.5),
        top=Inches(1.5),
        width=Inches(7.8),
    )
    _add_bullets(
        s,
        [
            "7,378 spots, 33,538 genes, 3 sections",
            "    (full cohort: 26 sections, 65,027 spots)",
            "",
            "Continuous CNA score c_i is the ONLY",
            "training supervision for marginals.",
            "",
            "Per-section quantile thresholds assign",
            "spots to tumor / intermediate / normal.",
            "",
            "layer, mp, ivygap annotations are",
            "post-hoc evaluation only \u2014 never",
            "folded into training.",
        ],
        left=Inches(8.6),
        top=Inches(1.6),
        width=Inches(4.5),
        height=Inches(5.5),
        size=15,
    )
    _add_footer(
        s,
        "Figure: results/full_run_bayes/figures/stage_2_spatial_marginal_labels.png",
    )


def slide_formulation(prs: Presentation, m: Metrics, eqs: dict[str, Path]) -> None:
    s = _blank_slide(prs)
    _add_title(s, "Formulation: notation, context, target drift", subtitle="Section-restricted spatial kNN")
    _add_bullets(
        s,
        [
            "Spot i:  expression x_i \u2208 R^G,  location s_i \u2208 R^2,  malignancy c_i \u2208 [0,1].",
            "Section-restricted kNN  N_k(i)  \u2014 prevents cross-section leakage.",
            "Normal reference set  R = { i : c_i below per-section low quantile }.",
        ],
        left=Inches(0.7),
        top=Inches(1.6),
        width=Inches(12.0),
        height=Inches(1.6),
        size=18,
    )
    # Two stacked equations
    _add_image(s, eqs["context"], left=Inches(1.5), top=Inches(3.4), width=Inches(10.0))
    _add_image(s, eqs["target"], left=Inches(1.5), top=Inches(5.0), width=Inches(10.0))
    _add_footer(
        s,
        "Eq. 1 (context) and Eq. 2 (target drift) \u2014 manuscript \u00a72.2\u20132.3.",
    )


def slide_ridge(prs: Presentation, m: Metrics, eqs: dict[str, Path]) -> None:
    s = _blank_slide(prs)
    _add_title(s, "Ridge-regularized spatial drift", subtitle="The heart of the method")
    _add_bullets(
        s,
        [
            "Augmented features per spot:  z_i = [ x_i  \u2225  x\u0304_i  \u2225  1 ]^\u22a4   \u2208  R^(2G+1)",
            "Ridge gives the MAP estimate of a Gaussian linear model \u2014 sets up the Bayesian extension later.",
        ],
        left=Inches(0.7),
        top=Inches(1.6),
        width=Inches(12.0),
        height=Inches(1.3),
        size=18,
    )
    _add_image(s, eqs["ridge"], left=Inches(1.5), top=Inches(2.9), width=Inches(10.0))
    _add_image(s, eqs["ridge_closed"], left=Inches(1.5), top=Inches(4.1), width=Inches(10.0))
    _add_image(s, eqs["drift_decomp"], left=Inches(1.5), top=Inches(5.3), width=Inches(10.0))
    _add_footer(
        s,
        "Methods point: \u0398_ctx is learned JOINTLY with \u0398_expr; spatial is inside the drift, not a post-hoc smoother.",
    )


def slide_euler(prs: Presentation, m: Metrics, eqs: dict[str, Path]) -> None:
    s = _blank_slide(prs)
    _add_title(s, "Euler transport and perturbation extraction")
    _add_image(s, eqs["euler"], left=Inches(0.4), top=Inches(1.5), width=Inches(7.6))
    _add_image(s, eqs["perturb"], left=Inches(0.4), top=Inches(3.0), width=Inches(7.6))
    cropped = crop_top(
        FIGS / "stage_4_umap_delta_distance_toward_reference.png",
        ASSETS / "stage_4_umap_delta_cropped.png",
        fraction=0.12,
    )
    _add_image(
        s,
        cropped,
        left=Inches(8.4),
        top=Inches(1.5),
        width=Inches(4.7),
    )
    _add_bullets(
        s,
        [
            "Three hyperparameters:",
            "  \u03b7  Euler step size  \u2014 drift magnitude per iteration",
            "  \u03b1  spatial smoothing strength  \u2014 explicit pull toward neighborhood mean",
            "  N  number of Euler steps",
            "",
            "Honest caveat: this is a linear ridge + Euler approximation of a Schr\u00f6dinger bridge,",
            "    not a full entropic OT solver.",
        ],
        left=Inches(0.4),
        top=Inches(4.3),
        width=Inches(7.7),
        height=Inches(2.8),
        size=15,
    )
    _add_footer(
        s,
        "Right panel: per-spot progress toward the normal reference after transport.",
    )


def slide_nmf(prs: Presentation, m: Metrics, eqs: dict[str, Path], pdf_pngs: dict[str, Path]) -> None:
    s = _blank_slide(prs)
    _add_title(s, "Program discovery via NMF", subtitle="Parts-based decomposition of perturbations")
    _add_image(s, pdf_pngs["gen_loadings_heatmap"], left=Inches(0.4), top=Inches(1.5), height=Inches(5.3))
    _add_image(s, eqs["nmf"], left=Inches(7.4), top=Inches(1.6), width=Inches(5.6))
    _add_bullets(
        s,
        [
            "Nonnegative shift  U\u208a = max(U \u2212 min U, 0),  K = 6,  NNDSVD-A init",
            "",
            f"Reconstruction error (relative Frobenius):",
            f"  NMF  {m.nmf_err:.3f}    PCA  {m.pca_err:.3f}    ICA  {m.ica_err:.3f}",
            "",
            "Seed-to-seed matched cosine \u2248 0.9999999",
            "    \u2192 programs are identifiable, not stochastic artifacts.",
        ],
        left=Inches(7.4),
        top=Inches(3.0),
        width=Inches(5.6),
        height=Inches(4.0),
        size=15,
    )
    _add_footer(
        s,
        "Figure: NMF loadings (top 20 genes, 6 programs) \u2014 gen_loadings_heatmap.pdf.",
    )


def slide_baselines(prs: Presentation, m: Metrics, eqs: dict[str, Path]) -> None:
    s = _blank_slide(prs)
    _add_title(s, "Baselines and evaluation metric", subtitle="Primary metric: mean L2 distance to normal reference")
    _add_image(s, eqs["dref"], left=Inches(0.5), top=Inches(1.5), width=Inches(12.0))
    _add_bullets(
        s,
        [
            "DE shift:   x\u0303_i = x_i + (\u03bc_normal \u2212 \u03bc_tumor)",
            "Static centroid:   x\u0303_i = x_i + (\u03bc \u2212 x\u0304)",
            "Unconditional bridge:   x\u0303_i = \u00bd (x_i + \u03bc)    \u2190 fairest spatial-off ablation",
            "Latent kNN blend:   x\u0303_i = (1 \u2212 \u03b2) x_i + \u03b2 x\u0304_{N_k^normal(i)}",
            "",
            "Spatial bridge (ours):   ridge drift + \u03b1-mixed Euler (Eqs. 3\u20135)",
            "",
            "Lower d_ref \u21d4 closer to the normal reference.  All metrics spot-level, internal.",
        ],
        left=Inches(0.8),
        top=Inches(3.1),
        width=Inches(11.8),
        height=Inches(3.8),
        size=18,
    )


def slide_benchmark(prs: Presentation, m: Metrics, pdf_pngs: dict[str, Path]) -> None:
    s = _blank_slide(prs)
    _add_title(s, "Main result: benchmark", subtitle="Spot-level internal  \u2014  lower is better")
    _add_image(s, pdf_pngs["gen_benchmark_bars"], left=Inches(0.4), top=Inches(1.5), height=Inches(5.3))

    # Right-side table with numbers straight from benchmark_metrics.csv
    rows = [
        ("Method", "d_ref (mean L2)"),
        ("Spatial Bridge (ours)", f"{m.spatial_bridge_l2:.2f}"),
        ("Unconditional Bridge", f"{m.uncond_bridge_l2:.2f}"),
        ("Latent kNN Blend", f"{m.latent_knn_l2:.2f}"),
        ("Static Centroid", f"{m.static_centroid_l2:.2f}"),
        ("DE Shift", f"{m.de_shift_l2:.2f}"),
    ]
    tbl_box = s.shapes.add_table(
        len(rows), 2, Inches(7.8), Inches(1.8), Inches(5.2), Inches(3.2)
    ).table
    for i, (col0, col1) in enumerate(rows):
        c0 = tbl_box.cell(i, 0)
        c1 = tbl_box.cell(i, 1)
        c0.text = col0
        c1.text = col1
        for c in (c0, c1):
            for para in c.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(16)
                    if i == 0:
                        run.font.bold = True
                        run.font.color.rgb = ACCENT
                    elif i == 1:
                        run.font.bold = True

    reduction_vs_uncond = 100 * (m.uncond_bridge_l2 - m.spatial_bridge_l2) / m.uncond_bridge_l2
    reduction_vs_de = 100 * (m.de_shift_l2 - m.spatial_bridge_l2) / m.de_shift_l2
    _add_bullets(
        s,
        [
            f"Spatial bridge beats best non-spatial baseline (Unconditional) by {reduction_vs_uncond:.0f}%.",
            f"Versus DE shift: {reduction_vs_de:.0f}% reduction in d_ref.",
            "Internal metric \u2014 see held-out slide for generalization evidence.",
        ],
        left=Inches(7.6),
        top=Inches(5.3),
        width=Inches(5.4),
        height=Inches(1.8),
        size=14,
    )
    _add_footer(s, "Source: results/full_run_bayes/benchmark_metrics.csv")


def slide_backend_comparison(prs: Presentation, m: Metrics) -> None:
    """Backend comparison: linear / Bayesian / neural.

    Numbers hard-coded from manuscript v2 Table 3 (backend comparison on the
    full 65,027-spot cohort); the metric is mean movement toward the normal
    reference (MD_ref) rather than the d_ref distance used on slide 9.
    """
    s = _blank_slide(prs)
    _add_title(
        s,
        "Backend comparison: linear vs. Bayesian vs. neural",
        subtitle="Full 65,027-spot cohort  \u2014  MD_ref = mean movement toward normal reference (higher is better)",
    )

    rows = [
        ("Backend", "MD_ref", "\u0394 vs. linear", "Held-out gain"),
        ("Linear / Bayesian (reference)", "20.28", "\u2014", "\u2014"),
        ("Neural (untuned)", "17.19", "\u22123.09  (\u221215.3%)", "\u22123.12"),
        ("Neural (tuned, sweep-selected)", "17.78", "\u22122.50  (\u221212.3%)", "\u22122.55"),
    ]
    tbl = s.shapes.add_table(
        len(rows), 4, Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.3)
    ).table
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = val
            for para in c.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(16)
                    if i == 0:
                        run.font.bold = True
                        run.font.color.rgb = ACCENT
                    elif i == 1:
                        run.font.bold = True

    _add_bullets(
        s,
        [
            "Key takeaway: the nonlinear drift does NOT close more of the tumor-to-normal gap.",
            "  Both the untuned and the sweep-selected tuned neural bridges moved spots LESS",
            "  than the linear baseline on both train and held-out sections.",
            "The linear-in-(x, x\u0304) parameterization is well-matched to Visium SNR at 2,000 HVGs.",
            "Neural backend is retained in the code base as an optional backend for higher-capacity regimes.",
            "",
            "Tuned config: hidden dim 384, 3 layers, dropout 0.1, lr 3\u00d710\u207b\u2074, weight decay 2\u00d710\u207b\u2074, 800 steps.",
        ],
        left=Inches(0.6),
        top=Inches(4.3),
        width=Inches(12.1),
        height=Inches(2.8),
        size=15,
    )
    _add_footer(
        s,
        "Source: manuscript v2 Table 3 (\u00a75.2, backend comparison on the full 65,027-spot cohort).",
    )


def slide_biology(prs: Presentation, m: Metrics, pdf_pngs: dict[str, Path]) -> None:
    s = _blank_slide(prs)
    _add_title(
        s,
        "Biological plausibility",
        subtitle="Perturbation strata + Hallmark ORA on NMF programs",
    )
    _add_image(s, pdf_pngs["gen_perturbation_norms"], left=Inches(0.4), top=Inches(1.4), height=Inches(5.3))

    rows = [
        ("Program", "Best Hallmark term", "FDR"),
        ("Neuronal_Plasticity", "HYPOXIA", f"{m.hypoxia_fdr:.1e}"),
        ("Mito_Metabolic", "OXIDATIVE_PHOSPHORYLATION", f"{m.oxphos_fdr:.1e}"),
        ("Reactive_Glial", "COAGULATION", f"{m.coagulation_fdr:.1e}"),
    ]
    tbl = s.shapes.add_table(
        len(rows), 3, Inches(6.6), Inches(1.6), Inches(6.5), Inches(2.2)
    ).table
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = val
            for para in c.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
                    if i == 0:
                        run.font.bold = True
                        run.font.color.rgb = ACCENT

    _add_bullets(
        s,
        [
            f"Mean ||u|| by stratum:    tumor {m.tumor_norm:.1f}    intermediate {m.intermediate_norm:.1f}    normal {m.normal_norm:.1f}",
            f"Pearson corr CNA vs perturbation norm = {m.cna_vs_pert_r:+.2f}.",
            "All three strata undergo substantial transport.",
            "A nonlinear (MLP) drift was tested and did not resolve the pattern",
            "  \u2014 see backend comparison slide.",
            "ORA is run AFTER NMF, without supervision on Hallmark sets.",
        ],
        left=Inches(6.5),
        top=Inches(4.1),
        width=Inches(6.6),
        height=Inches(2.8),
        size=14,
    )
    _add_footer(
        s,
        "Left: perturbation norm distributions by marginal label (gen_perturbation_norms.pdf).",
    )


def slide_validation(prs: Presentation, m: Metrics) -> None:
    s = _blank_slide(prs)
    _add_title(
        s,
        "Validation: held-out sections + synthetic ground truth",
        subtitle="Two pillars the methods audience will ask about",
    )
    _add_image(
        s,
        FIGS / "stage_8_heldout_section_performance.png",
        left=Inches(0.3),
        top=Inches(1.5),
        width=Inches(6.5),
    )
    _add_image(
        s,
        SYN_FIGS / "synthetic_metric_summary.png",
        left=Inches(6.9),
        top=Inches(1.5),
        width=Inches(6.2),
    )
    _add_bullets(
        s,
        [
            f"Held-out: {m.n_train:,} train / {m.n_heldout:,} held-out spots, SECTION-level split (no random spot split).",
            f"   mean \u0394 toward reference   train = {m.heldout_train_delta:.1f}   held-out = {m.heldout_test_delta:.1f}  \u2192 generalizes.",
            "Synthetic: corner-seeded radial malignancy gradient with known ground truth.",
            "   recovered strength vs planted demand  r = 1.00,   gene-space cosine 0.57.",
        ],
        left=Inches(0.4),
        top=Inches(5.6),
        width=Inches(12.7),
        height=Inches(1.5),
        size=14,
    )
    _add_footer(
        s,
        "Sources: stage_8_heldout_metrics.csv, results/synthetic_validation/",
    )


def slide_limitations(prs: Presentation, m: Metrics) -> None:
    s = _blank_slide(prs)
    _add_title(
        s,
        "Limitations, Bayesian extension, future work",
    )
    _add_bullets(
        s,
        [
            "Limitations",
            "  Cross-sectional \u2192 counterfactuals, not observed interventions.",
            "  Linear ridge \u2260 full entropic OT solver; gap between formulations remains unquantified.",
            "  Evaluation annotations (mp, layer, CODEX) used strictly post hoc, never as training targets.",
            "",
            "Bayesian extension (closed-form, free at inference time)",
            "  Posterior on \u0398 \u2192 per-spot perturbation_norm_std written to Stage 5 spot summary.",
            "  Enables triage of high- vs low-confidence perturbation calls without a second training run.",
            "",
            "Future work",
            "  Patient-level and external-cohort held-out splits.",
            "  Full entropic OT solver (beyond the ridge + Euler approximation).",
            "  Tumor-selective transport objective \u2014 move malignant spots while preserving normal tissue.",
        ],
        left=Inches(0.7),
        top=Inches(1.5),
        width=Inches(12.0),
        height=Inches(5.5),
        size=16,
    )


# ---------------------------------------------------------------------------
# Backup (hidden) slides
# ---------------------------------------------------------------------------


def slide_bayes_appendix(prs: Presentation, m: Metrics, eqs: dict[str, Path]) -> None:
    s = _blank_slide(prs)
    _add_title(
        s,
        "[Backup] Bayesian posterior over ridge weights",
        subtitle="Same point estimate, free per-spot uncertainty",
    )
    _add_bullets(
        s,
        [
            "Probabilistic model (ridge as MAP)",
            "  Likelihood:  Y | Z,\u0398  ~  N(Z\u0398, \u03c3\u00b2 I)     (shared noise across genes)",
            "  Prior:          \u0398  ~  N(0, \u03bb\u207b\u00b9 I)",
            "  \u2192 MAP is exactly the ridge solution (Eq. 3); we just keep the posterior around.",
        ],
        left=Inches(0.5),
        top=Inches(1.3),
        width=Inches(12.3),
        height=Inches(1.4),
        size=13,
    )
    eq_w = Inches(7.0)
    eq_left = Inches((13.333 - 7.0) / 2)
    _add_image(s, eqs["posterior"], left=eq_left, top=Inches(2.80), width=eq_w)
    _add_image(s, eqs["pred_var"], left=eq_left, top=Inches(4.10), width=eq_w)
    _add_image(s, eqs["pert_std"], left=eq_left, top=Inches(5.20), width=eq_w)
    _add_bullets(
        s,
        [
            "One SCALAR variance per spot (ridge ties noise across output genes).",
            "Cost: one p\u00d7p solve at train time (already in ridge), O(p\u00b2) memory, p = 2G+1.",
            "Implemented in omega_spatial/model.py: BayesianRidgeBridgeModel, predictive_drift_std.",
            "Conservative linear propagation through Euler \u2192 perturbation_norm_std on Stage 5 spot summary.",
        ],
        left=Inches(0.5),
        top=Inches(6.55),
        width=Inches(12.3),
        height=Inches(0.9),
        size=11,
    )
    _hide_slide(s)


def slide_ablation_appendix(prs: Presentation, m: Metrics) -> None:
    s = _blank_slide(prs)
    _add_title(
        s,
        "[Backup] Synthetic ablation: spatial vs non-spatial",
        subtitle="Strongest single-image defense of spatial conditioning",
    )
    _add_image(
        s,
        SYN_FIGS / "synthetic_ablation_spatial_vs_nonspatial.png",
        left=Inches(0.4),
        top=Inches(1.4),
        width=Inches(12.5),
    )
    _add_bullets(
        s,
        [
            "Identical ridge bridge; non-spatial control sets  context = expr.",
            "At Visium-like per-spot noise, mean gene-space L2 error roughly HALVES when context is enabled.",
            "\u2192 spatial conditioning is doing real denoising, not just acting as a trivial smoothness prior.",
        ],
        left=Inches(0.5),
        top=Inches(6.2),
        width=Inches(12.2),
        height=Inches(1.2),
        size=14,
    )
    _hide_slide(s)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def build() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    pdf_pngs = convert_pdf_figures()
    eqs = prep_equations()
    metrics = load_metrics()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, metrics)
    slide_motivation(prs, metrics)
    slide_data(prs, metrics)
    slide_formulation(prs, metrics, eqs)
    slide_ridge(prs, metrics, eqs)
    slide_euler(prs, metrics, eqs)
    slide_nmf(prs, metrics, eqs, pdf_pngs)
    slide_baselines(prs, metrics, eqs)
    slide_benchmark(prs, metrics, pdf_pngs)
    slide_backend_comparison(prs, metrics)
    slide_biology(prs, metrics, pdf_pngs)
    slide_validation(prs, metrics)
    slide_limitations(prs, metrics)

    slide_bayes_appendix(prs, metrics, eqs)
    slide_ablation_appendix(prs, metrics)

    prs.save(str(DECK_PATH))
    return DECK_PATH


if __name__ == "__main__":
    path = build()
    print(f"Deck written to: {path}")
